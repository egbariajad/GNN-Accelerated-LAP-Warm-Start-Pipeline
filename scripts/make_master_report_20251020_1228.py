#!/usr/bin/env python3
import argparse, glob, io
from datetime import datetime
from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

def parse_args():
    p = argparse.ArgumentParser(description="Auto-build a benchmarking PowerPoint from CSV logs.")
    p.add_argument("--inputs", nargs="+", required=True,
                   help="Glob patterns for CSV files, e.g. logs/performance/*.csv logs/experiments/*.csv")
    p.add_argument("--out", required=True, help="Output PPTX path")
    p.add_argument("--title", default="LA-H Bench – Auto Report")
    p.add_argument("--author", default="")
    p.add_argument("--dpi", type=int, default=170)
    return p.parse_args()

def load_all_csv(patterns):
    files = []
    for pat in patterns:
        files.extend(glob.glob(pat))
    frames = []
    for f in sorted(set(files)):
        try:
            df = pd.read_csv(f)
            if not isinstance(df, pd.DataFrame) or df.empty:
                continue
            df["__source_file"] = f
            frames.append(df)
        except Exception:
            continue
    if not frames:
        return pd.DataFrame()
    big = pd.concat(frames, ignore_index=True)
    big = normalize_columns(big)
    return big

def normalize_columns(df):
    out = df.copy()

    # drop duplicate column names early
    if out.columns.duplicated().any():
        out = out.loc[:, ~out.columns.duplicated()]

    rename_map = {
        "n":"size_n","size":"size_n","sizeN":"size_n","size_n":"size_n",
        "T_pred_ms":"T_pred_ms","pred_ms":"T_pred_ms",
        "T_xfer_ms":"T_xfer_ms","xfer_ms":"T_xfer_ms",
        "T_refine_ms":"T_refine_ms","refine_ms":"T_refine_ms",
        "T_total_ms":"T_total_ms","total_ms":"T_total_ms","seeded_ms":"T_total_ms",
        "scipylap_ms":"scipylap_ms","scipy_ms":"scipylap_ms",
        "T_H2D_ms":"T_H2D_ms","H2D_ms":"T_H2D_ms",
        "T_D2H_ms":"T_D2H_ms","D2H_ms":"T_D2H_ms",
        "BW_H2D_GBs":"BW_H2D_GBs","H2D_GBs":"BW_H2D_GBs","BW_H2D(GB/s)":"BW_H2D_GBs",
        "BW_D2H_GBs":"BW_D2H_GBs","D2H_GBs":"BW_D2H_GBs","BW_D2H(GB/s)":"BW_D2H_GBs",
        "bytes_H2D":"bytes_H2D","bytes_D2H":"bytes_D2H",
        "gpu":"gpu_name","gpu_name":"gpu_name","device":"gpu_name",
        "dtype":"dtype","precision":"dtype",
        "pinned":"pinned","pin":"pinned",
        "trial":"trial","run":"trial",
    }
    real_map = {c: rename_map[c] for c in out.columns if c in rename_map}
    out = out.rename(columns=real_map)

    # if renaming created duplicates, collapse again
    if out.columns.duplicated().any():
        out = out.loc[:, ~out.columns.duplicated()]

    def coerce_numeric_col(frame, col):
        obj = frame[col]
        if isinstance(obj, pd.DataFrame):
            # merge duplicated same-named columns by first non-null numeric
            merged = None
            for sub in obj.columns:
                s = pd.to_numeric(obj[sub], errors="coerce")
                merged = s if merged is None else merged.combine_first(s)
            frame[col] = merged
        else:
            frame[col] = pd.to_numeric(obj, errors="coerce")

    to_num = [
        "size_n","T_pred_ms","T_xfer_ms","T_refine_ms","T_total_ms","scipylap_ms",
        "T_H2D_ms","T_D2H_ms","BW_H2D_GBs","BW_D2H_GBs","bytes_H2D","bytes_D2H","trial"
    ]
    for c in to_num:
        if c in out.columns:
            coerce_numeric_col(out, c)

    for c in ["gpu_name","dtype","pinned","__source_file"]:
        if c in out.columns:
            out[c] = out[c].astype(str)

    if "size_n" in out.columns:
        out["size_n"] = pd.to_numeric(out["size_n"], errors="coerce").astype("Int64")

    return out

def add_speedup(df):
    if {"scipylap_ms","T_total_ms"}.issubset(df.columns):
        sp = df["scipylap_ms"] / df["T_total_ms"]
        df = df.assign(speedup_x=sp.replace([np.inf, -np.inf], np.nan))
    return df

def median_by_n(df, cols):
    df = df.dropna(subset=["size_n"])
    g = df.groupby("size_n", dropna=True)
    med = g[cols].median(numeric_only=True)
    cnt = g.size().rename("count")
    out = med.join(cnt, how="left").reset_index()
    return out.sort_values("size_n")

def fig_to_png_bytes(fig, dpi=170):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf

def plot_lines(x, ys, labels, title, ylabel):
    fig = plt.figure(figsize=(8.5, 4.8))
    for y, lab in zip(ys, labels):
        plt.plot(x, y, marker="o", label=lab)
    plt.title(title)
    plt.xlabel("n")
    plt.ylabel(ylabel)
    if len(ys) > 1:
        plt.legend()
    plt.grid(True, alpha=0.3)
    return fig

def table_image(df, title):
    fig = plt.figure(figsize=(9, 0.5 + 0.35*max(3, len(df))))
    plt.axis("off")
    plt.title(title, pad=10)
    tbl = plt.table(cellText=df.values, colLabels=df.columns, loc="center")
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(8)
    tbl.scale(1, 1.2)
    return fig

def add_title(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_picture_slide(prs, title, img_bytes, width_inches=9.0):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(img_bytes, Inches(0.5), Inches(1.2), width=Inches(width_inches-1.0))
    return slide

def add_two_pics_slide(prs, title, left_bytes, right_bytes):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(left_bytes, Inches(0.5), Inches(1.2), width=Inches(4.25))
    slide.shapes.add_picture(right_bytes, Inches(4.75), Inches(1.2), width=Inches(4.25))
    return slide

def section_txfer(prs, df, dpi):
    if "size_n" not in df.columns:
        return
    need_any = any(c in df.columns for c in ["T_H2D_ms","BW_H2D_GBs","T_D2H_ms","BW_D2H_GBs"])
    if not need_any:
        return

    cols = [c for c in ["T_H2D_ms","T_D2H_ms","BW_H2D_GBs","BW_D2H_GBs"] if c in df.columns]
    med = median_by_n(df, cols)

    img1 = img2 = None
    if "T_H2D_ms" in med.columns:
        fig = plot_lines(med["size_n"], [med["T_H2D_ms"]], ["H2D"], "Median H2D Transfer Time", "ms")
        img1 = fig_to_png_bytes(fig, dpi)
    if "T_D2H_ms" in med.columns:
        fig = plot_lines(med["size_n"], [med["T_D2H_ms"]], ["D2H"], "Median D2H Transfer Time", "ms")
        img2 = fig_to_png_bytes(fig, dpi)
    if img1 and img2:
        add_two_pics_slide(prs, "Transfers: Time (Median)", img1, img2)
    elif img1:
        add_picture_slide(prs, "Transfers: H2D Time (Median)", img1)
    elif img2:
        add_picture_slide(prs, "Transfers: D2H Time (Median)", img2)

    bw_imgs = []
    if "BW_H2D_GBs" in med.columns:
        fig = plot_lines(med["size_n"], [med["BW_H2D_GBs"]], ["H2D"], "Median H2D Throughput", "GB/s")
        bw_imgs.append(fig_to_png_bytes(fig, dpi))
    if "BW_D2H_GBs" in med.columns:
        fig = plot_lines(med["size_n"], [med["BW_D2H_GBs"]], ["D2H"], "Median D2H Throughput", "GB/s")
        bw_imgs.append(fig_to_png_bytes(fig, dpi))
    if len(bw_imgs) == 2:
        add_two_pics_slide(prs, "Transfers: Throughput (Median)", bw_imgs[0], bw_imgs[1])
    elif len(bw_imgs) == 1:
        add_picture_slide(prs, "Transfers: Throughput (Median)", bw_imgs[0])

    show_cols = [c for c in ["size_n","T_H2D_ms","BW_H2D_GBs","T_D2H_ms","BW_D2H_GBs","count"] if c in med.columns]
    if show_cols:
        fig = table_image(med[show_cols].round(3), "Transfers – Median Summary")
        add_picture_slide(prs, "Transfers – Median Summary", fig_to_png_bytes(fig, dpi))

def section_solver(prs, df, dpi):
    if "size_n" not in df.columns:
        return
    has_any = any(c in df.columns for c in ["scipylap_ms","T_total_ms","T_refine_ms","T_pred_ms","T_xfer_ms"])
    if not has_any:
        return

    df = add_speedup(df)
    cols = [c for c in ["scipylap_ms","T_total_ms","T_pred_ms","T_xfer_ms","T_refine_ms","speedup_x"] if c in df.columns]
    med = median_by_n(df, cols)

    y_cols = [c for c in ["scipylap_ms","T_total_ms"] if c in med.columns]
    if y_cols:
        fig = plot_lines(med["size_n"], [med[c] for c in y_cols], y_cols, "Solver Time (Median)", "ms")
        add_picture_slide(prs, "Solver Time (Median)", fig_to_png_bytes(fig, dpi))

    br_cols = [c for c in ["T_pred_ms","T_xfer_ms","T_refine_ms"] if c in med.columns]
    if br_cols:
        fig = plot_lines(med["size_n"], [med[c] for c in br_cols], br_cols, "Seeded Breakdown (Median)", "ms")
        add_picture_slide(prs, "Seeded Breakdown (Median)", fig_to_png_bytes(fig, dpi))

    if "speedup_x" in med.columns:
        fig = plot_lines(med["size_n"], [med["speedup_x"]], ["SciPy/Seeded"], "Speedup (Median)", "x")
        add_picture_slide(prs, "Speedup (Median)", fig_to_png_bytes(fig, dpi))

    show = ["size_n"] + y_cols + br_cols + (["speedup_x"] if "speedup_x" in med.columns else [])
    if show:
        fig = table_image(med[show].round(3), "Solver – Median Summary")
        add_picture_slide(prs, "Solver – Median Summary", fig_to_png_bytes(fig, dpi))

def section_meta(prs, df, dpi):
    rows = []
    if "dtype" in df.columns:
        rows.append(("Distinct dtypes", len(df["dtype"].dropna().unique())))
    if "gpu_name" in df.columns:
        rows.append(("Distinct GPUs", len(df["gpu_name"].dropna().unique())))
    if "pinned" in df.columns:
        rows.append(("Pinned flag variations", len(df["pinned"].dropna().unique())))
    if "size_n" in df.columns:
        rows.append(("Distinct n", len(df["size_n"].dropna().unique())))
    if not rows:
        return
    meta = pd.DataFrame(rows, columns=["Metric","Value"])
    fig = table_image(meta, "Dataset Coverage")
    add_picture_slide(prs, "Dataset Coverage", fig_to_png_bytes(fig, dpi))

def main():
    args = parse_args()
    df = load_all_csv(args.inputs)
    prs = Presentation()

    subtitle_bits = []
    if args.author:
        subtitle_bits.append(f"Author: {args.author}")
    subtitle_bits.append(f"Generated: {datetime.now().isoformat(timespec='seconds')}")
    add_title(prs, args.title, " | ".join(subtitle_bits))

    if df.empty or "size_n" not in df.columns:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "No data found"
    else:
        section_meta(prs, df, args.dpi)
        section_txfer(prs, df, args.dpi)
        section_solver(prs, df, args.dpi)

        tmp = (df.groupby("__source_file").size()
                 .rename("rows").reset_index()
                 .sort_values("rows", ascending=False))
        if not tmp.empty:
            fig = table_image(tmp, "Loaded CSV Files & Row Counts")
            add_picture_slide(prs, "Loaded CSV Files", fig_to_png_bytes(fig, args.dpi))

    out_path = Path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"Saved PowerPoint: {out_path}")

if __name__ == "__main__":
    main()
