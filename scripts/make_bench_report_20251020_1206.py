#!/usr/bin/env python3
import argparse, glob, io
from datetime import datetime
from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR

# ---------- Helpers ----------
def load_csvs(patterns):
    paths = []
    for pat in patterns:
        paths.extend(glob.glob(pat))
    out = []
    for p in sorted(set(paths)):
        try:
            df = pd.read_csv(p)
            if not df.empty:
                out.append((p, df))
        except Exception:
            pass
    return out

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_section(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide.shapes.title
    tx.text = title
    return slide

def add_text_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = line
        p.level = 0
    return slide

def df_to_table_slide(prs, title, df, max_rows=20):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rows = min(len(df), max_rows)
    cols = len(df.columns)
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    table = slide.shapes.add_table(rows+1, cols, left, top, width, height).table
    for j, col in enumerate(df.columns):
        table.cell(0, j).text = str(col)
    for i in range(rows):
        for j, col in enumerate(df.columns):
            val = df.iloc[i, j]
            table.cell(i+1, j).text = f"{val}"
    return slide

def add_image_slide(prs, title, image_bytes):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    pic = slide.shapes.add_picture(image_bytes, Inches(0.6), Inches(1.5), width=Inches(8.8))
    return slide

def fig_to_bytes():
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", dpi=180)
    plt.close()
    buf.seek(0)
    return buf

def safe_num(x):
    try:
        return float(x)
    except Exception:
        return np.nan

def has_cols(df, cols):
    return all(c in df.columns for c in cols)

def add_note(slide, text):
    shp = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(9), Inches(0.6))
    tf = shp.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.LEFT
    p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

# ---------- Detectors ----------
def detect_t_xfer(df):
    cand = [
        ("size_n","T_H2D_ms","BW_H2D_GBs","T_D2H_ms","BW_D2H_GBs"),
        ("n","T_H2D_ms","BW_H2D_GBs","T_D2H_ms","BW_D2H_GBs"),
        ("size","T_H2D_ms","BW_H2D_GBs","T_D2H_ms","BW_D2H_GBs"),
    ]
    for cols in cand:
        if has_cols(df, cols):
            return {"n":cols[0],"h2d":"T_H2D_ms","bh2d":"BW_H2D_GBs","d2h":"T_D2H_ms","bd2h":"BW_D2H_GBs"}
    # tolerate bandwidth-less schema
    cand2 = [("size_n","T_H2D_ms","T_D2H_ms"),("n","T_H2D_ms","T_D2H_ms")]
    for cols in cand2:
        if has_cols(df, cols):
            return {"n":cols[0],"h2d":"T_H2D_ms","bh2d":None,"d2h":"T_D2H_ms","bd2h":None}
    return None

def detect_baseline(df):
    # Typical columns from your earlier runs
    base_cols = ["size_n","T_pred_ms","T_xfer_ms","T_refine_ms","T_total_ms","scipylap_ms"]
    if has_cols(df, base_cols):
        return { "n":"size_n","tp":"T_pred_ms","tx":"T_xfer_ms","tr":"T_refine_ms","tt":"T_total_ms","sci":"scipylap_ms" }
    # Relaxed detection
    alt = [
        ("n","T_pred_ms","T_xfer_ms","T_refine_ms","T_total_ms","scipylap_ms"),
        ("size","T_pred_ms","T_xfer_ms","T_refine_ms","T_total_ms","scipylap_ms"),
    ]
    for cols in alt:
        if has_cols(df, cols):
            return { "n":cols[0],"tp":"T_pred_ms","tx":"T_xfer_ms","tr":"T_refine_ms","tt":"T_total_ms","sci":"scipylap_ms" }
    return None

def generic_time_columns(df):
    # Fallback: any *_ms columns vs size
    size_candidates = [c for c in df.columns if c in ("n","size","size_n")]
    time_cols = [c for c in df.columns if c.endswith("_ms")]
    return size_candidates[0] if size_candidates else None, time_cols

# ---------- Plotters ----------
def plot_txfer(df, m):
    g = df.copy()
    for k in [m["n"], m["h2d"], m["d2h"]]:
        if k and k in g.columns:
            g[k] = g[k].apply(safe_num)
    g = g.dropna(subset=[m["n"]])
    if m["h2d"] in g:
        g[m["h2d"]] = g[m["h2d"]].astype(float)
    if m["d2h"] in g:
        g[m["d2h"]] = g[m["d2h"]].astype(float)
    agg = g.groupby(m["n"]).agg(["median","mean","count"])
    plt.figure()
    if m["h2d"] in g:
        plt.plot(agg.index, agg[(m["h2d"],"median")], marker="o", label="H2D median (ms)")
    if m["d2h"] in g:
        plt.plot(agg.index, agg[(m["d2h"],"median")], marker="o", label="D2H median (ms)")
    plt.xlabel("n")
    plt.ylabel("Time (ms)")
    plt.title("Transfer Times by size (median)")
    plt.grid(True); plt.legend()
    img1 = fig_to_bytes()

    if m["bh2d"] or m["bd2h"]:
        plt.figure()
        if m["bh2d"] and m["bh2d"] in g:
            g[m["bh2d"]] = g[m["bh2d"]].astype(float)
            agg2 = g.groupby(m["n"])[m["bh2d"]].median()
            plt.plot(agg2.index, agg2.values, marker="s", label="H2D BW (GB/s)")
        if m["bd2h"] and m["bd2h"] in g:
            g[m["bd2h"]] = g[m["bd2h"]].astype(float)
            agg3 = g.groupby(m["n"])[m["bd2h"]].median()
            plt.plot(agg3.index, agg3.values, marker="s", label="D2H BW (GB/s)")
        plt.xlabel("n"); plt.ylabel("Bandwidth (GB/s)")
        plt.title("Transfer Bandwidth by size (median)")
        plt.grid(True); plt.legend()
        img2 = fig_to_bytes()
    else:
        img2 = None
    # summary table
    summary = []
    for n, sub in g.groupby(m["n"]):
        row = {"n": n}
        if m["h2d"] in g:
            row["T_H2D_med_ms"] = float(sub[m["h2d"]].median())
        if m["bh2d"] and m["bh2d"] in g:
            row["BW_H2D_med_GBs"] = float(sub[m["bh2d"]].median())
        if m["d2h"] in g:
            row["T_D2H_med_ms"] = float(sub[m["d2h"]].median())
        if m["bd2h"] and m["bd2h"] in g:
            row["BW_D2H_med_GBs"] = float(sub[m["bd2h"]].median())
        summary.append(row)
    summary_df = pd.DataFrame(summary).sort_values("n")
    return img1, img2, summary_df

def plot_baseline(df, m):
    g = df.copy()
    for k in m.values():
        if k and k in g:
            g[k] = g[k].apply(safe_num)
    g = g.dropna(subset=[m["n"]])
    agg = g.groupby(m["n"]).median(numeric_only=True)
    plt.figure()
    for key, label in [("tp","T_pred_ms"),("tx","T_xfer_ms"),("tr","T_refine_ms"),("tt","T_total_ms"),("sci","SciPy_ms")]:
        col = m.get(key)
        if col and col in agg:
            plt.plot(agg.index, agg[col], marker="o", label=label)
    plt.xlabel("n"); plt.ylabel("Median time (ms)")
    plt.title("Baseline Median Times by size")
    plt.grid(True); plt.legend()
    img = fig_to_bytes()

    # Speedup if SciPy and T_total exist
    sp = None
    if m.get("sci") and m.get("tt") and m["sci"] in agg and m["tt"] in agg:
        sp = (agg[m["sci"]] / agg[m["tt"]]).rename("speedup_vs_scipy")
        speed_df = pd.DataFrame({"n": agg.index, "speedup_vs_scipy": sp.values})
    else:
        speed_df = None
    return img, agg.reset_index(), speed_df

def plot_generic(df, size_col, time_cols):
    g = df.copy()
    g[size_col] = g[size_col].apply(safe_num)
    g = g.dropna(subset=[size_col])
    agg = g.groupby(size_col).median(numeric_only=True)
    cols = [c for c in time_cols if c in agg.columns]
    if not cols:
        return None, None
    plt.figure()
    for c in cols:
        plt.plot(agg.index, agg[c], marker="o", label=c)
    plt.xlabel(size_col); plt.ylabel("Median (ms)")
    plt.title("Timing medians by size")
    plt.grid(True); plt.legend()
    return fig_to_bytes(), agg.reset_index()

# ---------- Main ----------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--inputs", nargs="+", required=True, help="Glob(s) to CSVs")
    ap.add_argument("--out", type=str, default=None)
    args = ap.parse_args()

    ts = datetime.now().strftime("%Y%m%d_%H%M")
    out = args.out or f"artifacts/reports/bench_report_{ts}.pptx"
    Path("artifacts/reports").mkdir(parents=True, exist_ok=True)

    csvs = load_csvs(args.inputs)
    prs = Presentation()
    add_title_slide(prs, "GPU Benchmark Report",
                    f"Generated: {ts}\nInputs: {', '.join([p for p,_ in csvs]) or 'None'}")

    if not csvs:
        add_text_slide(prs, "No CSVs found", ["Nothing matched inputs."])
        prs.save(out); print("Saved:", out); return

    add_section(prs, "Transfer (H2D / D2H)")
    found_txfer = False
    for path, df in csvs:
        m = detect_t_xfer(df)
        if not m: continue
        found_txfer = True
        img1, img2, summ = plot_txfer(df, m)
        add_text_slide(prs, f"T_xfer: {Path(path).name}", [f"Rows: {len(df)}"])
        add_image_slide(prs, "Transfer times (median)", img1)
        if img2:
            add_image_slide(prs, "Transfer bandwidth (median)", img2)
        df_to_table_slide(prs, "Transfer summary (median by n)", summ)
    if not found_txfer:
        add_text_slide(prs, "No transfer CSVs detected", ["No matching columns for T_xfer."])

    add_section(prs, "Baseline Timings")
    found_base = False
    for path, df in csvs:
        m = detect_baseline(df)
        if not m: continue
        found_base = True
        img, agg_tbl, speed_tbl = plot_baseline(df, m)
        add_text_slide(prs, f"Baseline: {Path(path).name}", [f"Rows: {len(df)}"])
        add_image_slide(prs, "Median timings by size", img)
        df_to_table_slide(prs, "Median table", agg_tbl)
        if speed_tbl is not None:
            df_to_table_slide(prs, "Speedup vs SciPy (T_total)", speed_tbl)
    if not found_base:
        add_text_slide(prs, "No baseline CSVs detected", ["No matching columns for baseline schema."])

    add_section(prs, "Generic Timing CSVs")
    unknowns = []
    for path, df in csvs:
        if detect_t_xfer(df) or detect_baseline(df):
            continue
        size_col, tcols = generic_time_columns(df)
        if size_col and tcols:
            img, tbl = plot_generic(df, size_col, tcols)
            if img is not None:
                add_text_slide(prs, f"Generic: {Path(path).name}", [f"Rows: {len(df)}"])
                add_image_slide(prs, "Timing medians by size", img)
                df_to_table_slide(prs, "Median table", tbl.head(30))
            else:
                unknowns.append(path)
        else:
            unknowns.append(path)

    if unknowns:
        add_text_slide(prs, "Unknown CSV schemas (skipped)", unknowns[:30])

    prs.save(out)
    print("Saved PowerPoint:", out)

if __name__ == "__main__":
    main()
