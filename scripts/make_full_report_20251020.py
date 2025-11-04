#!/usr/bin/env python3
import argparse
import io
import glob
from datetime import datetime
from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --------- IO & parsing ---------
REQUIRED_BASELINE_KEYS = {"size_n","T_total_ms"}
OPTIONAL_BASELINE_KEYS = {"T_pred_ms","T_xfer_ms","T_refine_ms","scipylap_ms"}
TXFER_KEYS = {"n","T_H2D_ms","BW_H2D_GBs","T_D2H_ms","BW_D2H_GBs"}

def read_any_csv(path):
    try:
        df = pd.read_csv(path)
        df.columns = [c.strip() for c in df.columns]
        return df
    except Exception:
        return None

def classify_dataframe(df):
    cols = set(df.columns)
    if "size_n" in cols and "T_total_ms" in cols:
        return "baseline"
    if {"n","T_H2D_ms"}.issubset(cols):
        return "txfer"
    return "unknown"

def load_inputs(patterns):
    files = []
    for patt in patterns:
        files.extend(sorted(glob.glob(patt)))
    base, xfer, unknown = [], [], []
    for p in files:
        df = read_any_csv(p)
        if df is None or df.empty:
            continue
        kind = classify_dataframe(df)
        if kind == "baseline":
            base.append((p, df))
        elif kind == "txfer":
            xfer.append((p, df))
        else:
            unknown.append((p, df))
    return base, xfer, unknown

# --------- Aggregations ---------
def agg_baseline(df):
    keep = ["size_n","T_total_ms","T_pred_ms","T_xfer_ms","T_refine_ms","scipylap_ms"]
    df = df[[c for c in keep if c in df.columns]].copy()
    g = df.groupby("size_n")
    med = g.median(numeric_only=True).reset_index()
    # speedup if SciPy exists
    if "scipylap_ms" in med.columns:
        med["speedup_x"] = med["scipylap_ms"] / med["T_total_ms"]
    return med.sort_values("size_n")

def agg_txfer(df):
    keep = ["n","T_H2D_ms","BW_H2D_GBs","T_D2H_ms","BW_D2H_GBs"]
    df = df[[c for c in keep if c in df.columns]].copy()
    if "n" not in df.columns:
        return pd.DataFrame()
    g = df.groupby("n")
    med = g.median(numeric_only=True).reset_index()
    return med.sort_values("n")

# --------- Plot helpers ---------
def fig_to_png_bytes(fig, dpi=160):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf

def plot_stacked_breakdown(bmed):
    cols = [c for c in ["T_pred_ms","T_xfer_ms","T_refine_ms"] if c in bmed.columns]
    if not cols:
        return None
    fig, ax = plt.subplots(figsize=(7.5, 4.5))
    bottom = np.zeros(len(bmed))
    x = np.arange(len(bmed))
    for c in cols:
        ax.bar(x, bmed[c].values, bottom=bottom, label=c.replace("_ms",""))
        bottom += bmed[c].values
    ax.set_xticks(x, [str(n) for n in bmed["size_n"].values])
    ax.set_xlabel("n")
    ax.set_ylabel("Time (ms)")
    ax.set_title("Median breakdown per n")
    ax.legend()
    return fig_to_png_bytes(fig)

def plot_total_vs_scipy(bmed):
    if "scipylap_ms" not in bmed.columns:
        return None
    fig, ax = plt.subplots(figsize=(7.5, 4.5))
    ax.plot(bmed["size_n"], bmed["T_total_ms"], marker="o", label="T_total")
    ax.plot(bmed["size_n"], bmed["scipylap_ms"], marker="o", label="SciPy")
    ax.set_xlabel("n"); ax.set_ylabel("ms"); ax.set_title("Total vs SciPy (median)")
    ax.legend()
    return fig_to_png_bytes(fig)

def plot_speedup(bmed):
    if "speedup_x" not in bmed.columns:
        return None
    fig, ax = plt.subplots(figsize=(7.5, 4.0))
    ax.plot(bmed["size_n"], bmed["speedup_x"], marker="o")
    ax.axhline(1.0, linestyle="--")
    ax.set_xlabel("n"); ax.set_ylabel("×"); ax.set_title("Speedup (SciPy / T_total)")
    return fig_to_png_bytes(fig)

def plot_txfer_bw(tmed):
    if tmed.empty:
        return None
    fig, ax = plt.subplots(figsize=(7.5, 4.5))
    ax.plot(tmed["n"], tmed.get("BW_H2D_GBs", pd.Series([np.nan]*len(tmed))), marker="o", label="H2D GB/s")
    ax.plot(tmed["n"], tmed.get("BW_D2H_GBs", pd.Series([np.nan]*len(tmed))), marker="o", label="D2H GB/s")
    ax.set_xlabel("n"); ax.set_ylabel("GB/s"); ax.set_title("Median transfer bandwidth")
    ax.legend()
    return fig_to_png_bytes(fig)

def plot_txfer_time(tmed):
    if tmed.empty:
        return None
    fig, ax = plt.subplots(figsize=(7.5, 4.5))
    ax.plot(tmed["n"], tmed.get("T_H2D_ms", pd.Series([np.nan]*len(tmed))), marker="o", label="T_H2D ms")
    ax.plot(tmed["n"], tmed.get("T_D2H_ms", pd.Series([np.nan]*len(tmed))), marker="o", label="T_D2H ms")
    ax.set_xlabel("n"); ax.set_ylabel("ms"); ax.set_title("Median transfer time")
    ax.legend()
    return fig_to_png_bytes(fig)

# --------- PowerPoint helpers ---------
def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_text_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = line
        p.level = 0

def add_table_slide(prs, title, df, max_rows=14):
    chunks = [df.iloc[i:i+max_rows] for i in range(0, len(df), max_rows)]
    for idx, part in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{title}" + (f" (part {idx+1})" if len(chunks) > 1 else "")
        rows, cols = part.shape
        left, top, width, height = Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.0)
        table = slide.shapes.add_table(rows+1, cols, left, top, width, height).table
        table.columns[0].width = Inches(1.2)
        # headers
        for j, col in enumerate(part.columns):
            table.cell(0, j).text = str(col)
        # rows
        for i in range(rows):
            for j in range(cols):
                val = part.iloc[i, j]
                if isinstance(val, float):
                    table.cell(i+1, j).text = f"{val:.3f}"
                else:
                    table.cell(i+1, j).text = str(val)

def add_image_slide(prs, title, png_bytes):
    if png_bytes is None:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(png_bytes, Inches(0.6), Inches(1.2), width=Inches(8.8))

# --------- Main ---------
def main():
    ap = argparse.ArgumentParser(description="Auto-generate full benchmark PowerPoint from CSV logs.")
    ap.add_argument("--inputs", nargs="+", required=True, help="CSV globs, e.g. logs/performance/*.csv")
    ap.add_argument("--out", type=str, default=None, help="Output pptx path")
    args = ap.parse_args()

    base, xfer, unknown = load_inputs(args.inputs)
    now = datetime.now().strftime("%Y%m%d_%H%M")
    out = args.out or f"artifacts/reports/bench_report_{now}.pptx"
    Path(out).parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    add_title_slide(
        prs,
        "GPU Benchmark Report",
        f"Generated: {now}\nInputs: {', '.join([p for p,_ in base+xfer])}"
    )

    # Baseline aggregation
    if base:
        add_text_slide(prs, "Baseline files", [p for p,_ in base])

        # Concatenate and aggregate
        bcat = pd.concat([df for _, df in base], ignore_index=True)
        bcat = bcat.replace([np.inf, -np.inf], np.nan).dropna(how="all", axis=1)
        bmed = agg_baseline(bcat)
        if not bmed.empty:
            add_table_slide(prs, "Baseline (median per n)", bmed)
            add_image_slide(prs, "Breakdown: T_pred/T_xfer/T_refine", plot_stacked_breakdown(bmed))
            add_image_slide(prs, "Total vs SciPy", plot_total_vs_scipy(bmed))
            add_image_slide(prs, "Speedup (SciPy / T_total)", plot_speedup(bmed))
        else:
            add_text_slide(prs, "Baseline summary", ["No aggregatable columns found."])
    else:
        add_text_slide(prs, "Baseline", ["No baseline CSVs detected."])

    # Transfer aggregation
    if xfer:
        add_text_slide(prs, "Transfer files", [p for p,_ in xfer])

        xcat = pd.concat([df for _, df in xfer], ignore_index=True)
        xcat = xcat.replace([np.inf, -np.inf], np.nan).dropna(how="all", axis=1)
        xmed = agg_txfer(xcat)
        if not xmed.empty:
            add_table_slide(prs, "Transfer (median per n)", xmed)
            add_image_slide(prs, "Transfer bandwidth", plot_txfer_bw(xmed))
            add_image_slide(prs, "Transfer time", plot_txfer_time(xmed))
        else:
            add_text_slide(prs, "Transfer summary", ["No aggregatable columns found."])
    else:
        add_text_slide(prs, "Transfer", ["No transfer CSVs detected."])

    # Unknown files note
    if unknown:
        add_text_slide(prs, "Unknown CSV schemas (skipped)", [p for p,_ in unknown])

    prs.save(out)
    print(f"Saved PowerPoint: {out}")

if __name__ == "__main__":
    main()
