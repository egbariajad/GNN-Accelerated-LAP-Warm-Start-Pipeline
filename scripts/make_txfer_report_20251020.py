#!/usr/bin/env python3
import argparse
import glob
import io
from datetime import datetime
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def parse_args():
    p = argparse.ArgumentParser(description="Build a PowerPoint report from H2D/D2H transfer benchmark CSVs.")
    p.add_argument("--inputs", nargs="*", default=["logs/performance/t_xfer_*.csv"],
                   help="Glob patterns to CSVs (default: logs/performance/t_xfer_*.csv)")
    p.add_argument("--outdir", type=Path, default=Path("artifacts/reports"),
                   help="Output directory for the .pptx (default: artifacts/reports)")
    p.add_argument("--title", type=str, default="GPU Transfer Benchmark Report",
                   help="Deck title")
    p.add_argument("--author", type=str, default="Auto-Report",
                   help="Author shown on title slide")
    p.add_argument("--dpi", type=int, default=150, help="Figure DPI for embedded charts")
    return p.parse_args()


def load_all(patterns):
    frames = []
    for pat in patterns:
        for fp in glob.glob(pat):
            df = pd.read_csv(fp)
            df["__source__"] = fp
            frames.append(df)
    if not frames:
        raise FileNotFoundError("No CSV files matched the provided patterns.")
    df = pd.concat(frames, ignore_index=True)

    if "pinned" in df.columns:
        df["pinned"] = (
            df["pinned"]
            .astype(str)
            .replace({"1": "pinned", "0": "unpinned", "True": "pinned", "False": "unpinned"})
        )
    else:
        df["pinned"] = "unknown"

    if "dtype" not in df.columns:
        df["dtype"] = "fp32"

    for col in ["size_n", "bytes"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce")
    for col in ["T_H2D_ms", "T_D2H_ms", "BW_H2D_GBps", "BW_D2H_GBps"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce")

    df = df.dropna(subset=["size_n", "T_H2D_ms", "T_D2H_ms"])
    df["size_n"] = df["size_n"].astype(int)
    return df


def median_summary(df):
    keys = ["size_n", "dtype", "pinned"]
    agg = {
        "T_H2D_ms": "median",
        "BW_H2D_GBps": "median",
        "T_D2H_ms": "median",
        "BW_D2H_GBps": "median",
        "bytes": "median",
    }
    g = df.groupby(keys, dropna=False).agg(agg).reset_index()
    g = g.rename(columns={
        "T_H2D_ms": "T_H2D_med_ms",
        "BW_H2D_GBps": "BW_H2D_med_GBps",
        "T_D2H_ms": "T_D2H_med_ms",
        "BW_D2H_GBps": "BW_D2H_med_GBps",
        "bytes": "bytes_med",
    })
    return g.sort_values(["dtype", "pinned", "size_n"]).reset_index(drop=True)


def _fig_bytes():
    return io.BytesIO()


def line_chart(df, x, y, hue_cols, title, dpi=150):
    fig = plt.figure()
    hue = df[hue_cols].astype(str).agg("|".join, axis=1)
    for key, sub in df.groupby(hue):
        sub = sub.sort_values(x)
        plt.plot(sub[x], sub[y], marker="o", label=key)
    plt.title(title)
    plt.xlabel(x)
    plt.ylabel(y)
    plt.legend()
    plt.grid(True, linestyle="--", linewidth=0.5, alpha=0.6)
    bio = _fig_bytes()
    fig.savefig(bio, format="png", dpi=dpi, bbox_inches="tight")
    plt.close(fig)
    bio.seek(0)
    return bio


def bar_chart(df, x, y, title, dpi=150):
    fig = plt.figure()
    cats = sorted(df[x].unique())
    xpos = np.arange(len(cats))
    plt.bar(xpos, df.set_index(x).reindex(cats)[y].values, width=0.7)
    plt.title(title)
    plt.xlabel(x)
    plt.ylabel(y)
    plt.xticks(xpos, [str(c) for c in cats])
    plt.grid(True, axis="y", linestyle="--", linewidth=0.5, alpha=0.6)
    bio = _fig_bytes()
    fig.savefig(bio, format="png", dpi=dpi, bbox_inches="tight")
    plt.close(fig)
    bio.seek(0)
    return bio


def make_table_data(summary):
    cols = ["size_n", "dtype", "pinned", "T_H2D_med_ms", "BW_H2D_med_GBps", "T_D2H_med_ms", "BW_D2H_med_GBps"]
    tbl = summary[cols].copy()
    tbl["T_H2D_med_ms"] = tbl["T_H2D_med_ms"].map(lambda v: f"{v:.3f}")
    tbl["BW_H2D_med_GBps"] = tbl["BW_H2D_med_GBps"].map(lambda v: f"{v:.2f}")
    tbl["T_D2H_med_ms"] = tbl["T_D2H_med_ms"].map(lambda v: f"{v:.3f}")
    tbl["BW_D2H_med_GBps"] = tbl["BW_D2H_med_GBps"].map(lambda v: f"{v:.2f}")
    return tbl


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0


def add_picture_slide(prs, title, image_bytes):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    left = Inches(0.6)
    top = Inches(1.5)
    slide.shapes.add_picture(image_bytes, left, top, width=Inches(9.2))


def add_table_slide(prs, title, df):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9.5), Inches(5.0)).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.text_frame.paragraphs[0].font.bold = True
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = str(df.iat[i, j])
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.text_frame.paragraphs[0].font.size = Pt(12)


def main():
    args = parse_args()
    df = load_all(args.inputs)
    summary = median_summary(df)

    now = datetime.now().strftime("%Y%m%d_%H%M")
    args.outdir.mkdir(parents=True, exist_ok=True)
    out_pptx = args.outdir / f"t_xfer_report_{now}.pptx"

    prs = Presentation()
    add_title_slide(prs, args.title, f"Generated: {now}\nAuthor: {args.author}")

    sources = sorted(df["__source__"].unique().tolist())
    device_names = sorted(df.get("device", pd.Series(["unknown"])).fillna("unknown").unique().tolist())
    bullets = [f"Files: {len(sources)}"]
    bullets += sources[:10]
    bullets.append(f"Device(s): {', '.join(device_names)}")
    add_bullet_slide(prs, "Data Sources & Context", bullets)

    tbl = make_table_data(summary)
    max_rows = 18
    if len(tbl) <= max_rows:
        add_table_slide(prs, "Median Summary (All)", tbl)
    else:
        chunks = int(np.ceil(len(tbl) / max_rows))
        for i in range(chunks):
            sub = tbl.iloc[i * max_rows : (i + 1) * max_rows]
            add_table_slide(prs, f"Median Summary (Part {i+1}/{chunks})", sub)

    hue_cols = ["dtype", "pinned"]

    img = line_chart(summary, "size_n", "T_H2D_med_ms", hue_cols, "Median H2D Time vs Size (ms)", dpi=args.dpi)
    add_picture_slide(prs, "H2D Time (Median)", img)

    img = line_chart(summary, "size_n", "T_D2H_med_ms", hue_cols, "Median D2H Time vs Size (ms)", dpi=args.dpi)
    add_picture_slide(prs, "D2H Time (Median)", img)

    img = line_chart(summary, "size_n", "BW_H2D_med_GBps", hue_cols, "Median H2D Bandwidth (GB/s)", dpi=args.dpi)
    add_picture_slide(prs, "H2D Bandwidth (Median)", img)

    img = line_chart(summary, "size_n", "BW_D2H_med_GBps", hue_cols, "Median D2H Bandwidth (GB/s)", dpi=args.dpi)
    add_picture_slide(prs, "D2H Bandwidth (Median)", img)

    for (dtype, pinned), sub in summary.groupby(["dtype", "pinned"], dropna=False):
        ttl = f"Per-Config Summary: dtype={dtype}, pinned={pinned}"
        img = bar_chart(sub, "size_n", "T_H2D_med_ms", ttl + " (H2D ms)", dpi=args.dpi)
        add_picture_slide(prs, ttl + " (H2D Time)", img)
        img = bar_chart(sub, "size_n", "T_D2H_med_ms", ttl + " (D2H ms)", dpi=args.dpi)
        add_picture_slide(prs, ttl + " (D2H Time)", img)

    prs.save(out_pptx)
    print(f"Saved PowerPoint: {out_pptx}")


if __name__ == "__main__":
    main()
