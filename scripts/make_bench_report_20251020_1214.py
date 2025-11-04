#!/usr/bin/env python3
import argparse, glob, io
from datetime import datetime
from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def load_many_csv(patterns):
    files = []
    for pat in patterns:
        files.extend(glob.glob(pat))
    if not files:
        raise FileNotFoundError(f"No CSVs matched: {patterns}")
    frames = []
    for fp in sorted(files):
        try:
            df = pd.read_csv(fp)
            if df is None or df.empty:
                continue
            df["__source_file"] = Path(fp).name
            frames.append(df)
        except Exception as e:
            print(f"[WARN] failed reading {fp}: {e}")
    if not frames:
        raise RuntimeError("No CSVs could be read")
    big = pd.concat(frames, ignore_index=True)
    big.columns = [c.strip().replace(" ", "_") for c in big.columns]
    return big

def coerce_numeric(df):
    out = df.copy()
    for c in out.columns:
        if c.startswith("__"):
            continue
        if out[c].dtype.kind in "biufc":
            continue
        if any(k in c.lower() for k in ["_ms","bw_","bytes_","mb","gb","speedup","size","n","trial"]):
            try:
                out[c] = pd.to_numeric(out[c], errors="coerce")
            except Exception:
                pass
    return out

def find_size_col(df):
    for cand in ["size_n","n","size","N","Size","matrix_n"]:
        if cand in df.columns:
            return cand
    raise KeyError("Couldn't find a size column (e.g., size_n or n).")

def numeric_groupby(df, by):
    num = df.select_dtypes(include=np.number).copy()
    if by not in num.columns and by in df.columns:
        num[by] = df[by]
    num = num.dropna(subset=[by])
    grouped = num.groupby(by).agg(["median","mean","count"])
    grouped.columns = [f"{a}_{b}" for a,b in grouped.columns]
    grouped = grouped.reset_index()
    return grouped

def fig_to_png_bytes(fig, dpi=160):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=dpi)
    plt.close(fig)
    buf.seek(0)
    return buf.read()

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_picture_slide(prs, title, png_bytes, width_inches=10.0):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(0.5), Inches(1.3), width=Inches(width_inches))
    return slide

def add_table_slide(prs, title, df, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rows, cols = df.shape
    left, top, width, height = Inches(0.5), Inches(1.3), Inches(10), Inches(0.8 + 0.3*rows)
    table = slide.shapes.add_table(rows+1, cols, left, top, width, height).table
    for j, col in enumerate(df.columns):
        table.cell(0, j).text = str(col)
    for i in range(rows):
        for j in range(cols):
            table.cell(i+1, j).text = str(df.iloc[i, j])
    if note:
        tx = slide.shapes.add_textbox(Inches(0.5), top+height+Inches(0.2), Inches(10), Inches(1.0))
        p = tx.text_frame.paragraphs[0]
        p.text = note
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT

def plot_txfer(df, size_col):
    cols_h2d = [c for c in df.columns if c.lower() in ["t_h2d_ms","th2d_ms","h2d_ms"]]
    cols_d2h = [c for c in df.columns if c.lower() in ["t_d2h_ms","td2h_ms","d2h_ms"]]
    bw_h2d   = [c for c in df.columns if "bw_h2d" in c.lower()]
    bw_d2h   = [c for c in df.columns if "bw_d2h" in c.lower()]
    if not cols_h2d and not bw_h2d and not cols_d2h and not bw_d2h:
        return None, None, None
    g = coerce_numeric(df).dropna(subset=[size_col])
    agg = numeric_groupby(g, size_col)
    img1 = img2 = None
    if cols_h2d:
        c = cols_h2d[0]
        fig = plt.figure(figsize=(8,4.5))
        plt.plot(agg[size_col], agg.get(f"{c}_median", np.nan), marker="o")
        plt.title("Host-to-Device transfer time (median)")
        plt.xlabel("n"); plt.ylabel("ms"); plt.grid(True, alpha=0.3)
        img1 = fig_to_png_bytes(fig)
    elif bw_h2d:
        c = bw_h2d[0]
        fig = plt.figure(figsize=(8,4.5))
        plt.plot(agg[size_col], agg.get(f"{c}_median", np.nan), marker="o")
        plt.title("Host-to-Device bandwidth (median)")
        plt.xlabel("n"); plt.ylabel("GB/s"); plt.grid(True, alpha=0.3)
        img1 = fig_to_png_bytes(fig)
    if cols_d2h:
        c = cols_d2h[0]
        fig = plt.figure(figsize=(8,4.5))
        plt.plot(agg[size_col], agg.get(f"{c}_median", np.nan), marker="o")
        plt.title("Device-to-Host transfer time (median)")
        plt.xlabel("n"); plt.ylabel("ms"); plt.grid(True, alpha=0.3)
        img2 = fig_to_png_bytes(fig)
    elif bw_d2h:
        c = bw_d2h[0]
        fig = plt.figure(figsize=(8,4.5))
        plt.plot(agg[size_col], agg.get(f"{c}_median", np.nan), marker="o")
        plt.title("Device-to-Host bandwidth (median)")
        plt.xlabel("n"); plt.ylabel("GB/s"); plt.grid(True, alpha=0.3)
        img2 = fig_to_png_bytes(fig)
    keep = [size_col]
    for c in cols_h2d + cols_d2h + bw_h2d + bw_d2h:
        mc = f"{c}_median"
        if mc in agg.columns:
            keep.append(mc)
    summ = agg[keep] if len(keep) > 1 else agg[[size_col]]
    return img1, img2, summ

def plot_runtime(df, size_col):
    cand = ["t_pred_ms","t_xfer_ms","t_refine_ms","t_total_ms","scipylap_ms"]
    present = [c for c in cand if c in df.columns]
    if not present:
        return (None, None, None)
    g = coerce_numeric(df).dropna(subset=[size_col])
    agg = numeric_groupby(g, size_col)
    fig = plt.figure(figsize=(9,5))
    for c in present:
        mc = f"{c}_median"
        if mc in agg.columns:
            plt.plot(agg[size_col], agg[mc], marker="o", label=c.replace("_ms",""))
    plt.title("Runtime components (median)")
    plt.xlabel("n"); plt.ylabel("ms"); plt.grid(True, alpha=0.3); plt.legend()
    img = fig_to_png_bytes(fig)
    img2, tbl = None, None
    if "scipylap_ms_median" in agg.columns and "t_total_ms_median" in agg.columns:
        agg["speedup_x"] = agg["scipylap_ms_median"] / agg["t_total_ms_median"]
        fig2 = plt.figure(figsize=(8,4.5))
        plt.plot(agg[size_col], agg["speedup_x"], marker="o")
        plt.title("Speedup vs SciPy (median)")
        plt.xlabel("n"); plt.ylabel("x"); plt.grid(True, alpha=0.3)
        img2 = fig_to_png_bytes(fig2)
        tbl = agg[[size_col, "scipylap_ms_median","t_total_ms_median","speedup_x"]]
    return (img, img2, tbl)

def main():
    p = argparse.ArgumentParser("Auto bench-to-PowerPoint")
    p.add_argument("--inputs", nargs="+", required=True, help="Glob patterns for CSVs")
    p.add_argument("--out", required=True, help="Output .pptx")
    p.add_argument("--title", default="Benchmark Report")
    p.add_argument("--subtitle", default="")
    p.add_argument("--dpi", type=int, default=160)
    args = p.parse_args()

    df = load_many_csv(args.inputs)
    df = coerce_numeric(df)
    size_col = find_size_col(df)

    prs = Presentation()
    ts = datetime.now().strftime("%Y-%m-%d %H:%M")
    add_title_slide(prs, args.title, args.subtitle or f"Generated {ts}")

    coverage = df[[size_col]].copy()
    coverage["count"] = 1
    cov = coverage.groupby(size_col)["count"].sum().reset_index()
    add_table_slide(prs, "Dataset coverage", cov.rename(columns={"count":"rows"}))

    t1, t2, t_tbl = plot_txfer(df, size_col)
    if t1 is not None:
        add_picture_slide(prs, "Transfers – H2D", t1)
    if t2 is not None:
        add_picture_slide(prs, "Transfers – D2H", t2)
    if t_tbl is not None and not t_tbl.empty:
        add_table_slide(prs, "Transfers – median summary", t_tbl)

    r_img, sp_img, r_tbl = plot_runtime(df, size_col)
    if r_img is not None:
        add_picture_slide(prs, "Runtime components", r_img)
    if sp_img is not None:
        add_picture_slide(prs, "Speedup vs SciPy", sp_img)
    if r_tbl is not None:
        add_table_slide(prs, "Runtime summary (median)", r_tbl)

    Path(args.out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(f"Saved PowerPoint: {args.out}")

if __name__ == "__main__":
    main()
